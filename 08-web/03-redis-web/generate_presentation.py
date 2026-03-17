#!/usr/bin/env python3
"""Generate a visually rich PPTX presentation for the Redis Log System (Light Theme)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Light color palette ───────────────────────────────────────────────────
SLIDE_BG     = RGBColor(0xF8, 0xF9, 0xFC)  # very light blue-gray
TITLE_BG     = RGBColor(0x2C, 0x3E, 0x50)  # dark navy for title slide
ACCENT_BLUE  = RGBColor(0x1A, 0x73, 0xE8)  # Google blue
ACCENT_TEAL  = RGBColor(0x00, 0x9B, 0x8D)  # deep teal
ACCENT_ORANGE= RGBColor(0xE8, 0x71, 0x0A)  # warm orange
ACCENT_RED   = RGBColor(0xDC, 0x36, 0x2E)  # red
ACCENT_PURPLE= RGBColor(0x7B, 0x4F, 0xBF)  # purple
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_PRIMARY  = RGBColor(0x1F, 0x2A, 0x37)  # near-black
TEXT_SECONDARY= RGBColor(0x4B, 0x55, 0x63)  # dark gray
TEXT_MUTED    = RGBColor(0x6B, 0x72, 0x80)  # medium gray
BORDER_COLOR  = RGBColor(0xE0, 0xE4, 0xEB)  # light border
CARD_BG      = RGBColor(0xFF, 0xFF, 0xFF)   # white cards
CARD_BG_ALT  = RGBColor(0xF1, 0xF5, 0xF9)  # slightly tinted cards
CODE_BG      = RGBColor(0x1E, 0x29, 0x3B)   # dark code blocks for contrast
ACCENT_BAR   = RGBColor(0xE2, 0xE8, 0xF0)   # subtle accent bar bg

# Soft tinted backgrounds for architecture cards
BLUE_TINT    = RGBColor(0xEB, 0xF5, 0xFF)
TEAL_TINT    = RGBColor(0xE6, 0xFA, 0xF5)
ORANGE_TINT  = RGBColor(0xFF, 0xF4, 0xE6)
RED_TINT     = RGBColor(0xFE, 0xF2, 0xF2)
PURPLE_TINT  = RGBColor(0xF3, 0xEE, 0xFC)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helper functions ───────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, radius=None, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def set_text(shape, text, size=18, color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return tf

def add_text_box(slide, left, top, width, height, text, size=18, color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return tf

def add_paragraph(tf, text, size=18, color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI', space_before=Pt(6)):
    p = tf.add_paragraph()
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return p

def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    shape = add_shape(slide, left, top, width, height, CODE_BG, radius=0.02)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)

    lines = code_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(0)

        keywords = ['#include', 'namespace', 'class', 'void', 'auto', 'const',
                     'try', 'catch', 'return', 'while', 'if', 'for', 'bool',
                     'int', 'size_t', 'std::', 'constexpr', 'using', 'public:',
                     'private:', 'struct', 'true', 'false', 'break', 'else',
                     'explicit', 'template', 'typename']

        run = p.add_run()
        run.text = line if line else ' '
        run.font.size = Pt(font_size)
        run.font.name = 'Consolas'

        has_keyword = any(kw in line for kw in keywords)
        has_string = '"' in line or "'" in line
        has_comment = line.strip().startswith('//')

        if has_comment:
            run.font.color.rgb = RGBColor(0x6A, 0x99, 0x55)
        elif has_keyword:
            run.font.color.rgb = RGBColor(0x7D, 0xCF, 0xFF)
        elif has_string:
            run.font.color.rgb = RGBColor(0xCE, 0x91, 0x78)
        else:
            run.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)

    return shape

def add_arrow_right(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_section_title(slide, text, subtitle=None):
    set_slide_bg(slide, SLIDE_BG)
    # Decorative accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.10), H, ACCENT_BLUE)
    # Title
    add_text_box(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
                 text, size=32, color=TEXT_PRIMARY, bold=True)
    # Thin line under title
    add_rect(slide, Inches(0.6), Inches(1.05), Inches(3), Inches(0.04), ACCENT_TEAL)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(1.15), Inches(12), Inches(0.45),
                     subtitle, size=16, color=TEXT_MUTED)

def add_bullet_card(slide, left, top, width, height, title, bullets, title_color=ACCENT_BLUE, card_color=CARD_BG, tint_color=None):
    shape = add_shape(slide, left, top, width, height, card_color, radius=0.03, border_color=BORDER_COLOR)
    if tint_color:
        shape.fill.fore_color.rgb = tint_color
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = title_color
    run.font.name = 'Segoe UI'

    for bullet in bullets:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"  {bullet}"
        run.font.size = Pt(13)
        run.font.color.rgb = TEXT_SECONDARY
        run.font.name = 'Segoe UI'

    return shape


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, WHITE)

# Top accent bar
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)
# Bottom accent bar
add_rect(slide, Inches(0), Inches(7.44), W, Inches(0.06), ACCENT_TEAL)

# Soft background shape
add_rect(slide, Inches(0), Inches(0.06), W, Inches(3.2), RGBColor(0xF0, 0xF7, 0xFF))

# Decorative circles (soft)
add_circle(slide, Inches(10.5), Inches(0.5), Inches(2.5), RGBColor(0xE8, 0xF0, 0xFE))
add_circle(slide, Inches(11.2), Inches(1.0), Inches(1.5), RGBColor(0xDB, 0xEA, 0xFE))

# Title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Redis Log System", size=52, color=TEXT_PRIMARY, bold=True)

# Subtitle
add_text_box(slide, Inches(1), Inches(3.1), Inches(11), Inches(0.6),
             "Клиент-серверная система сбора и просмотра логов", size=24, color=ACCENT_BLUE)

# Description
add_text_box(slide, Inches(1), Inches(4.2), Inches(8), Inches(1.5),
             "C++ · Redis Streams · SSE · Modern Web UI", size=18, color=TEXT_MUTED)

# Tech badges
badges = [("C++23", ACCENT_BLUE), ("Redis", ACCENT_RED), ("httplib", ACCENT_ORANGE),
          ("SSE", ACCENT_TEAL), ("HTML/CSS/JS", ACCENT_PURPLE)]
for i, (label, color) in enumerate(badges):
    bx = Inches(1 + i * 2.1)
    shape = add_shape(slide, bx, Inches(5.4), Inches(1.9), Inches(0.55), color, radius=0.15)
    set_text(shape, label, size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — ОБЗОР СИСТЕМЫ
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Обзор системы", "Что делает Redis Log System и зачем она нужна")

# Left card
add_bullet_card(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.8),
                "Проблема",
                ["Множество устройств генерируют логи",
                 "Нужен централизованный сбор в реальном времени",
                 "Просмотр и скачивание логов через веб-интерфейс",
                 "Масштабируемость и автоматическая очистка данных"],
                title_color=ACCENT_ORANGE, tint_color=ORANGE_TINT)

# Right card
add_bullet_card(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.8),
                "Решение",
                ["Redis Streams как шина сообщений",
                 "C++ клиент отправляет логи в Redis",
                 "C++ HTTP-сервер с SSE для real-time доставки",
                 "Современный веб-UI для просмотра и фильтрации"],
                title_color=ACCENT_TEAL, tint_color=TEAL_TINT)

# Bottom — flow
flow_y = Inches(5.0)
items = [
    ("Устройство\n(C++ клиент)", ACCENT_BLUE),
    ("Redis\nStreams", ACCENT_RED),
    ("HTTP-сервер\n(C++ / httplib)", ACCENT_ORANGE),
    ("Браузер\n(SSE + JS)", ACCENT_TEAL),
]
for i, (label, color) in enumerate(items):
    bx = Inches(0.8 + i * 3.2)
    shape = add_shape(slide, bx, flow_y, Inches(2.2), Inches(1.0), color, radius=0.05)
    set_text(shape, label, size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    if i < len(items) - 1:
        add_arrow_right(slide, Inches(3.1 + i * 3.2), Inches(5.22), Inches(0.7), Inches(0.5), BORDER_COLOR)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — АРХИТЕКТУРА
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Архитектура", "Как компоненты взаимодействуют друг с другом")

# --- Architecture boxes ---
box_h = Inches(3.5)
box_y = Inches(2.2)

# Client box
shape = add_shape(slide, Inches(0.4), box_y, Inches(3.0), box_h, BLUE_TINT, radius=0.03, border_color=ACCENT_BLUE)
tf = set_text(shape, "CLIENT", size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "", size=6, color=TEXT_PRIMARY)
add_paragraph(tf, "  41-redis-log-client", size=12, color=TEXT_PRIMARY, bold=True)
add_paragraph(tf, "", size=4, color=TEXT_PRIMARY)
add_paragraph(tf, "  • Подключение к Redis", size=11, color=TEXT_SECONDARY)
add_paragraph(tf, "  • Регистрация устройства", size=11, color=TEXT_SECONDARY)
add_paragraph(tf, "  • Генерация stream-ключа", size=11, color=TEXT_SECONDARY)
add_paragraph(tf, "  • XADD в Redis Stream", size=11, color=TEXT_SECONDARY)
add_paragraph(tf, "  • Интерактивный ввод", size=11, color=TEXT_SECONDARY)

# Arrows client -> Redis
add_arrow_right(slide, Inches(3.5), Inches(3.0), Inches(0.8), Inches(0.4), ACCENT_BLUE)
add_text_box(slide, Inches(3.5), Inches(2.55), Inches(0.9), Inches(0.35),
             "XADD", size=10, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_arrow_right(slide, Inches(3.5), Inches(4.3), Inches(0.8), Inches(0.4), ACCENT_TEAL)
add_text_box(slide, Inches(3.5), Inches(4.7), Inches(0.9), Inches(0.35),
             "HSET\nSADD", size=10, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

# Redis box
shape = add_shape(slide, Inches(4.4), box_y, Inches(3.8), box_h, RED_TINT, radius=0.03, border_color=ACCENT_RED)
tf = set_text(shape, "REDIS", size=14, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "", size=6, color=TEXT_PRIMARY)
add_paragraph(tf, "  Структуры данных:", size=12, color=TEXT_PRIMARY, bold=True)
add_paragraph(tf, "", size=4, color=TEXT_PRIMARY)
add_paragraph(tf, "  HASH    all_devices", size=11, color=ACCENT_ORANGE, font_name='Consolas')
add_paragraph(tf, "  SET     log_names:{dev}", size=11, color=ACCENT_TEAL, font_name='Consolas')
add_paragraph(tf, "  STREAM  log:{dev}:{file}", size=11, color=ACCENT_BLUE, font_name='Consolas')
add_paragraph(tf, "", size=4, color=TEXT_PRIMARY)
add_paragraph(tf, "  Keyspace Notifications", size=11, color=ACCENT_PURPLE)
add_paragraph(tf, "  TTL: 72 часа", size=11, color=TEXT_MUTED)

# Arrows Redis -> Server
add_arrow_right(slide, Inches(8.3), Inches(3.0), Inches(0.8), Inches(0.4), ACCENT_ORANGE)
add_text_box(slide, Inches(8.2), Inches(2.55), Inches(1.0), Inches(0.35),
             "XRANGE\nXREAD", size=10, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_arrow_right(slide, Inches(8.3), Inches(4.3), Inches(0.8), Inches(0.4), ACCENT_PURPLE)
add_text_box(slide, Inches(8.2), Inches(4.7), Inches(1.0), Inches(0.35),
             "SUBSCRIBE\npsubscribe", size=10, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# Server box
shape = add_shape(slide, Inches(9.2), box_y, Inches(3.8), box_h, ORANGE_TINT, radius=0.03, border_color=ACCENT_ORANGE)
tf = set_text(shape, "WEB SERVER", size=14, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "", size=6, color=TEXT_PRIMARY)
add_paragraph(tf, "  03-redis-web", size=12, color=TEXT_PRIMARY, bold=True)
add_paragraph(tf, "", size=4, color=TEXT_PRIMARY)
add_paragraph(tf, "  • notification_center", size=11, color=TEXT_SECONDARY)
add_paragraph(tf, "  • unified_stream_provider", size=11, color=TEXT_SECONDARY)
add_paragraph(tf, "  • unified_stream_handler", size=11, color=TEXT_SECONDARY)
add_paragraph(tf, "  • logs_download_handler", size=11, color=TEXT_SECONDARY)
add_paragraph(tf, "  • Статика: HTML/CSS/JS", size=11, color=TEXT_SECONDARY)

# Bottom: Browser
shape = add_shape(slide, Inches(9.7), Inches(6.0), Inches(2.8), Inches(0.85), TEAL_TINT, radius=0.05, border_color=ACCENT_TEAL)
tf = set_text(shape, "  BROWSER  (EventSource / SSE)", size=12, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

# Arrow server -> browser
connector_shape = slide.shapes.add_shape(
    MSO_SHAPE.DOWN_ARROW, Inches(10.9), Inches(5.65), Inches(0.35), Inches(0.35)
)
connector_shape.fill.solid()
connector_shape.fill.fore_color.rgb = ACCENT_TEAL
connector_shape.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — REDIS DATA MODEL
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Модель данных Redis", "Три типа ключей для хранения логов")

cards_data = [
    ("HASH", "all_devices", ACCENT_ORANGE, ORANGE_TINT,
     "Реестр всех устройств",
     [("Ключ", "all_devices"),
      ("Формат", 'device_name → platform'),
      ("Пример", '"anastasia_device" → "MacOS"'),
      ("TTL", "72 часа")]),
    ("SET", "log_names:{device}", ACCENT_TEAL, TEAL_TINT,
     "Список лог-файлов устройства",
     [("Ключ", "log_names:anastasia_device"),
      ("Элементы", "Имена лог-стримов"),
      ("Пример", '"log_17_03_2026T14_30_45.txt"'),
      ("TTL", "72 часа")]),
    ("STREAM", "log:{device}:{file}", ACCENT_BLUE, BLUE_TINT,
     "Содержимое лог-файла",
     [("Ключ", "log:anastasia_device:log_17_03..."),
      ("Поле", 'message'),
      ("Пример", '"[14:30:45.123] INFO Started"'),
      ("TTL", "72 часа")]),
]

for i, (redis_type, key_pattern, color, tint, description, fields) in enumerate(cards_data):
    cx = Inches(0.5 + i * 4.2)
    cy = Inches(1.7)

    # Type badge
    badge = add_shape(slide, cx, cy, Inches(1.3), Inches(0.45), color, radius=0.15)
    set_text(badge, redis_type, size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Card
    shape = add_shape(slide, cx, Inches(2.3), Inches(3.8), Inches(4.2), tint, radius=0.03, border_color=BORDER_COLOR)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = key_pattern
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = 'Consolas'

    add_paragraph(tf, description, size=12, color=TEXT_MUTED, space_before=Pt(10))
    add_paragraph(tf, "", size=6, color=TEXT_PRIMARY)

    for field_name, field_val in fields:
        p = tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = f"{field_name}:  "
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = TEXT_SECONDARY
        run.font.name = 'Segoe UI'
        run = p.add_run()
        run.text = field_val
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_PRIMARY
        run.font.name = 'Consolas'


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — КЛИЕНТ: ОБЗОР
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Клиент: 41-redis-log-client", "Отправка логов в Redis Stream")

# Left: explanation cards
cards = [
    ("1. Подключение", "Читает строку подключения из\nlog_client_config.txt", ACCENT_BLUE),
    ("2. Регистрация", "HSET all_devices + SADD log_names\nсо сгенерированным именем стрима", ACCENT_TEAL),
    ("3. Логирование", "XADD в Redis Stream с таймстампом\nи уровнем (INFO/WARN/ERROR/DEBUG)", ACCENT_ORANGE),
    ("4. Интерактив", "Цикл чтения stdin — пользователь\nвводит сообщения в реальном времени", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(cards):
    cy = Inches(1.65 + i * 1.4)
    # Color dot
    add_circle(slide, Inches(0.6), cy + Inches(0.1), Inches(0.35), color)
    add_text_box(slide, Inches(0.6), cy + Inches(0.1), Inches(0.35), Inches(0.35),
                 str(i+1), size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Text
    add_text_box(slide, Inches(1.15), cy, Inches(4.5), Inches(0.35),
                 title, size=16, color=color, bold=True)
    add_text_box(slide, Inches(1.15), cy + Inches(0.35), Inches(4.5), Inches(0.7),
                 desc, size=12, color=TEXT_SECONDARY)

# Right: code
code = '''namespace om {

constexpr auto key_ttl = std::chrono::hours(72);

void log_to_redis(redis::Redis& client,
                  const std::string& stream_key,
                  const std::string& level,
                  const std::string& message)
{
    auto now = std::chrono::system_clock::now();
    auto tm  = to_local_tm(now);
    auto ms  = std::chrono::duration_cast<
        std::chrono::milliseconds>(
            now.time_since_epoch()) % 1000;

    auto formatted = std::format(
        "[{:02}:{:02}:{:02}.{:03}] {} {}",
        tm.tm_hour, tm.tm_min, tm.tm_sec,
        ms.count(), level, message);

    client.xadd(stream_key, "*",
        {std::make_pair("message", formatted)});
}

} // namespace om'''

add_code_block(slide, Inches(6.0), Inches(1.65), Inches(6.8), Inches(5.3), code, font_size=10)
add_text_box(slide, Inches(6.0), Inches(1.3), Inches(3), Inches(0.35),
             "main.cxx — ключевая функция", size=11, color=TEXT_MUTED)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — КЛИЕНТ: FLOW
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Клиент: поток выполнения", "От запуска до завершения")

code_main = '''int main() {
    // 1. Читаем конфигурацию
    std::ifstream config_file("log_client_config.txt");
    std::getline(config_file, connection_string);

    // 2. Подключаемся к Redis
    redis::Redis redis_client(connection_string);

    // 3. Генерируем уникальное имя лог-стрима
    std::string log_file_name = om::generate_log_file_name();
    // "log_17_03_2026T14_30_45.txt"

    std::string stream_key =
        "log:" + device_name + ":" + log_file_name;

    // 4. Регистрируем устройство
    redis_client.hset("all_devices", device_name, platform);
    redis_client.sadd("log_names:" + device_name, log_file_name);

    // 5. Отправляем начальные логи
    om::log_to_redis(client, stream_key, "INFO", "Started");
    om::log_to_redis(client, stream_key, "DEBUG", "Init...");
    om::log_to_redis(client, stream_key, "WARNING", "No cfg");
    om::log_to_redis(client, stream_key, "ERROR", "DB fail!");

    // 6. Интерактивный режим
    while (std::getline(std::cin, user_input)) {
        if (user_input == "quit") break;
        om::log_to_redis(client, stream_key, "INFO", user_input);
    }
}'''
add_code_block(slide, Inches(0.6), Inches(1.65), Inches(8.0), Inches(5.3), code_main, font_size=10)

# Right side: key points
tints = [BLUE_TINT, TEAL_TINT, ORANGE_TINT, RED_TINT]
items_right = [
    ("Конфигурация", "Строка подключения из файла:\ntcp://default:pass@host:6379", ACCENT_BLUE),
    ("Имя стрима", "Автогенерация по timestamp:\nlog_DD_MM_YYYYTHH_MM_SS.txt", ACCENT_TEAL),
    ("Метаданные", "HASH: реестр устройств\nSET: список лог-файлов", ACCENT_ORANGE),
    ("TTL", "72 часа на все ключи\nАвтоматическая очистка", ACCENT_RED),
]
for i, (title, desc, color) in enumerate(items_right):
    cy = Inches(1.8 + i * 1.35)
    shape = add_shape(slide, Inches(9.0), cy, Inches(3.8), Inches(1.15), tints[i], radius=0.03, border_color=BORDER_COLOR)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = 'Segoe UI'
    add_paragraph(tf, desc, size=11, color=TEXT_SECONDARY, space_before=Pt(4))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — СЕРВЕР: NOTIFICATION CENTER
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Сервер: notification_center", "Подписка на Keyspace Notifications Redis")

# Left: explanation
add_bullet_card(slide, Inches(0.6), Inches(1.7), Inches(5.5), Inches(5.2),
                "Как работает notification_center",
                ["Redis Keyspace Notifications (KEA) уведомляют",
                 "  об изменениях в ключах all_devices и log_names:*",
                 "",
                 "Подписка через PSUBSCRIBE на паттерны:",
                 "  __keyspace@0__:all_devices",
                 "  __keyspace@0__:log_names:*",
                 "",
                 "При изменении — обновляет version-счётчики",
                 "  и будит SSE-потоки через condition_variable",
                 "",
                 "Работает в отдельном std::jthread",
                 "  с поддержкой std::stop_token",
                 "",
                 "Автопереподключение при ошибках Redis"],
                title_color=ACCENT_PURPLE, tint_color=PURPLE_TINT)

# Right: code
code_nc = '''class notification_center {
    std::mutex              mtx;
    std::condition_variable cv_devices;
    std::condition_variable cv_streams;
    std::condition_variable cv_any;

    std::atomic<uint64_t> devices_version{0};
    std::atomic<uint64_t> any_version{0};

    std::jthread worker_thread;

    void start() {
        redis_client->command(
            "CONFIG", "SET",
            "notify-keyspace-events", "KEA");

        worker_thread = std::jthread(
            [this](std::stop_token stoken) {
                auto sub = redis_client->subscriber();
                sub.on_pmessage([this](
                    std::string pattern,
                    std::string channel,
                    std::string msg) {
                    // Обновляем версии, будим CV
                    devices_version++;
                    cv_devices.notify_all();
                });
                sub.psubscribe("__keyspace@0__:*");
                while (!stoken.stop_requested())
                    sub.consume();
            });
    }
};'''
add_code_block(slide, Inches(6.4), Inches(1.7), Inches(6.4), Inches(5.2), code_nc, font_size=10)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — СЕРВЕР: SSE STREAMING
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Сервер: SSE-стриминг", "Единый эндпоинт /api/stream для всех данных")

# Top explanation
add_text_box(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.4),
             "unified_stream_provider — мультиплексирует devices, log_names и logs в одном SSE-соединении",
             size=14, color=TEXT_SECONDARY)

# Event cards
events = [
    ("devices_init\ndevices_update", "Список устройств\nHGETALL all_devices", ACCENT_BLUE, BLUE_TINT),
    ("log_names_init\nlog_names_update", "Имена логов устройства\nSMEMBERS log_names:{dev}", ACCENT_TEAL, TEAL_TINT),
    ("logs_init", "Начальная загрузка логов\nXRANGE пакетами по 1000", ACCENT_ORANGE, ORANGE_TINT),
    ("logs_new", "Новые записи в реальном времени\nXREAD с блокирующим таймаутом", ACCENT_PURPLE, PURPLE_TINT),
]

for i, (event_name, desc, color, tint) in enumerate(events):
    cx = Inches(0.4 + i * 3.2)
    cy = Inches(2.3)
    shape = add_shape(slide, cx, cy, Inches(3.0), Inches(1.8), tint, radius=0.03, border_color=BORDER_COLOR)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = event_name
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = 'Consolas'
    add_paragraph(tf, desc, size=11, color=TEXT_SECONDARY, space_before=Pt(8))

# Code
code_sse = '''bool operator()(size_t offset, httplib::DataSink& sink) const {
    // 1. Начальные данные
    send_sse(sink, "devices_init", get_devices_json());

    if (!device_.empty())
        send_sse(sink, "log_names_init", get_log_names_json());

    // 2. Пакетная загрузка существующих логов (XRANGE)
    while (true) {
        std::vector<stream_item> batch;
        redis_client_->xrange(stream_key_, start, "+",
                              batch_size_, std::back_inserter(batch));
        if (batch.empty()) break;
        send_sse(sink, first ? "logs_init" : "logs_new",
                 items_to_json(batch), last_log_id);
    }

    // 3. Real-time polling (XREAD + блокировка)
    while (true) {
        redis_client_->xread(stream_key_, last_log_id,
                             poll_interval_, std::back_inserter(data));
        if (!data.empty())
            send_sse(sink, "logs_new", items_to_json(items), last_id);
        // + обновления devices и log_names через notification_center
    }
}'''
add_code_block(slide, Inches(0.6), Inches(4.4), Inches(12.2), Inches(2.8), code_sse, font_size=10)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — ПОЧЕМУ ОДИН SSE-ЭНДПОИНТ
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Дизайн: единый SSE-эндпоинт", "Почему один /api/stream вместо нескольких")

# Problem vs Solution
add_bullet_card(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.5),
                "Проблема: лимит браузера",
                ["Браузеры ограничивают ~6 TCP-соединений на домен",
                 "EventSource (SSE) держит соединение открытым",
                 "3 отдельных SSE = 3 из 6 слотов заняты",
                 "Остальные запросы (CSS, JS, API) блокируются"],
                title_color=ACCENT_RED, tint_color=RED_TINT)

add_bullet_card(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.5),
                "Решение: мультиплексирование",
                ["Один SSE-эндпоинт для всех типов данных",
                 "event: devices_init / log_names_init / logs_new",
                 "JS-клиент разделяет данные по типу события",
                 "Всего 1 TCP-соединение вместо 3"],
                title_color=ACCENT_TEAL, tint_color=TEAL_TINT)

# Visual comparison
add_text_box(slide, Inches(0.6), Inches(4.6), Inches(5.5), Inches(0.4),
             "БЕЗ мультиплексирования:", size=14, color=ACCENT_RED, bold=True)
for i, label in enumerate(["GET /api/devices/stream", "GET /api/log_names/stream", "GET /api/logs/stream"]):
    shape = add_shape(slide, Inches(0.6), Inches(5.1 + i * 0.55), Inches(5.5), Inches(0.45),
                      RED_TINT, radius=0.03, border_color=ACCENT_RED)
    set_text(shape, f"  TCP #{i+1}  →  {label}", size=11, color=ACCENT_RED, font_name='Consolas')

add_text_box(slide, Inches(6.8), Inches(4.6), Inches(5.5), Inches(0.4),
             "С мультиплексированием:", size=14, color=ACCENT_TEAL, bold=True)
shape = add_shape(slide, Inches(6.8), Inches(5.1), Inches(5.8), Inches(1.55),
                  TEAL_TINT, radius=0.03, border_color=ACCENT_TEAL)
tf = set_text(shape, "  TCP #1  →  GET /api/stream?device=X&log_name=Y", size=11, color=ACCENT_TEAL, font_name='Consolas')
add_paragraph(tf, "", size=6, color=TEXT_PRIMARY)
add_paragraph(tf, "    event: devices_init  |  event: log_names_update", size=11, color=TEXT_SECONDARY, font_name='Consolas')
add_paragraph(tf, "    event: logs_init     |  event: logs_new", size=11, color=TEXT_SECONDARY, font_name='Consolas')


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — WEB UI
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Web-интерфейс", "Три панели: устройства, лог-файлы, содержимое")

# Simulate UI panels
ui_y = Inches(1.7)
ui_h = Inches(4.8)

# Shadow/border for whole UI
add_shape(slide, Inches(0.55), ui_y - Inches(0.02), Inches(12.3), ui_h + Inches(0.04),
          RGBColor(0xE8, 0xEC, 0xF1), radius=0.01, border_color=BORDER_COLOR)

# Header bar
add_rect(slide, Inches(0.6), ui_y, Inches(12.2), Inches(0.5), RGBColor(0x2C, 0x3E, 0x50))
add_text_box(slide, Inches(0.8), ui_y + Inches(0.05), Inches(4), Inches(0.4),
             "Device Logs Viewer", size=14, color=WHITE, bold=True)

# Panel 1: Devices
p1_x = Inches(0.6)
p1_w = Inches(2.5)
shape = add_shape(slide, p1_x, ui_y + Inches(0.5), p1_w, ui_h - Inches(0.5), RGBColor(0xFF, 0xFF, 0xFF), radius=0.01)
shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
shape.line.width = Pt(1)

add_text_box(slide, p1_x + Inches(0.1), ui_y + Inches(0.6), p1_w - Inches(0.2), Inches(0.3),
             "Devices", size=12, color=RGBColor(0x44, 0x44, 0x44), bold=True)
# Search box
search = add_shape(slide, p1_x + Inches(0.1), ui_y + Inches(1.0), p1_w - Inches(0.2), Inches(0.35),
                   RGBColor(0xFA, 0xFA, 0xFA), radius=0.1)
search.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
set_text(search, "  Search devices...", size=10, color=RGBColor(0xAA, 0xAA, 0xAA))
# Devices list
devices_list = ["anastasia_device", "ivan_device", "server_prod_01"]
for j, d in enumerate(devices_list):
    dy = ui_y + Inches(1.55 + j * 0.4)
    is_selected = j == 0
    bg_c = RGBColor(0xE3, 0xF2, 0xFD) if is_selected else RGBColor(0xFF, 0xFF, 0xFF)
    item = add_rect(slide, p1_x + Inches(0.05), dy, p1_w - Inches(0.1), Inches(0.35), bg_c)
    if is_selected:
        add_rect(slide, p1_x + Inches(0.05), dy, Inches(0.05), Inches(0.35), RGBColor(0x21, 0x96, 0xF3))
    add_text_box(slide, p1_x + Inches(0.15), dy, p1_w - Inches(0.3), Inches(0.35),
                 d, size=10, color=RGBColor(0x33, 0x33, 0x33), bold=is_selected)

# Panel 2: Log files
p2_x = Inches(3.2)
p2_w = Inches(3.0)
shape = add_shape(slide, p2_x, ui_y + Inches(0.5), p2_w, ui_h - Inches(0.5), RGBColor(0xFF, 0xFF, 0xFF), radius=0.01)
shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
shape.line.width = Pt(1)

add_text_box(slide, p2_x + Inches(0.1), ui_y + Inches(0.6), p2_w - Inches(0.2), Inches(0.3),
             "Log files", size=12, color=RGBColor(0x44, 0x44, 0x44), bold=True)

logs_list = ["log_17_03_2026T14_30_45.txt", "log_16_03_2026T09_15_22.txt", "log_15_03_2026T20_45_01.txt"]
for j, l in enumerate(logs_list):
    ly = ui_y + Inches(1.05 + j * 0.4)
    is_selected = j == 0
    bg_c = RGBColor(0xE3, 0xF2, 0xFD) if is_selected else RGBColor(0xFF, 0xFF, 0xFF)
    add_rect(slide, p2_x + Inches(0.05), ly, p2_w - Inches(0.1), Inches(0.35), bg_c)
    if is_selected:
        add_rect(slide, p2_x + Inches(0.05), ly, Inches(0.05), Inches(0.35), RGBColor(0x21, 0x96, 0xF3))
    add_text_box(slide, p2_x + Inches(0.15), ly, p2_w - Inches(0.3), Inches(0.35),
                 l, size=9, color=RGBColor(0x33, 0x33, 0x33), bold=is_selected)

# Panel 3: Log output (keep dark — it's a terminal!)
p3_x = Inches(6.3)
p3_w = Inches(6.5)
shape = add_shape(slide, p3_x, ui_y + Inches(0.5), p3_w, ui_h - Inches(0.5), RGBColor(0x1E, 0x1E, 0x1E), radius=0.01)

add_text_box(slide, p3_x + Inches(0.1), ui_y + Inches(0.55), Inches(1.5), Inches(0.3),
             "Log Output", size=12, color=RGBColor(0xAA, 0xAA, 0xAA), bold=True)

log_lines = [
    ("[14:30:45.123] INFO Application started successfully.", RGBColor(0x4F, 0xC1, 0xFF)),
    ("[14:30:45.135] DEBUG Initializing internal modules...", RGBColor(0x85, 0x85, 0x85)),
    ("[14:30:45.150] WARNING Config file not found", RGBColor(0xCC, 0xA7, 0x00)),
    ("[14:30:45.170] ERROR Failed to connect to secondary DB!", RGBColor(0xF4, 0x87, 0x71)),
    ("[14:30:46.001] INFO User logged in from 192.168.1.42", RGBColor(0x4F, 0xC1, 0xFF)),
    ("[14:30:46.050] INFO Processing batch #1247...", RGBColor(0x4F, 0xC1, 0xFF)),
    ("[14:30:46.200] DEBUG Memory usage: 142MB / 512MB", RGBColor(0x85, 0x85, 0x85)),
    ("[14:30:47.010] WARNING Slow query detected (450ms)", RGBColor(0xCC, 0xA7, 0x00)),
]
for j, (line, color) in enumerate(log_lines):
    add_text_box(slide, p3_x + Inches(0.15), ui_y + Inches(1.05 + j * 0.38), p3_w - Inches(0.3), Inches(0.35),
                 line, size=10, color=color, font_name='Consolas')

# Labels below
add_text_box(slide, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             "Поиск устройств  ·  Фильтр по платформе  ·  Auto-scroll  ·  Скачивание логов  ·  Подсветка уровней  ·  Анимация новых записей",
             size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — JS CLIENT
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Фронтенд: app.js", "EventSource и обработка SSE-событий")

code_js = '''// Единое SSE-соединение
function connectToStream() {
    let url = '/api/stream';
    const params = new URLSearchParams();
    if (currentDevice) params.set('device', currentDevice);
    if (currentStream) params.set('log_name', currentStream);
    if (params.toString()) url += '?' + params;

    eventSource = new EventSource(url);

    // Обработка устройств
    eventSource.addEventListener('devices_init', (e) => {
        allDevicesData = JSON.parse(e.data);
        renderPlatformFilters(platforms);
        filterAndRenderDevices();
    });

    // Обработка списка лог-файлов
    eventSource.addEventListener('log_names_init', (e) => {
        const names = JSON.parse(e.data);
        renderList(streamsList, names.sort().reverse());
    });

    // Начальная загрузка логов
    eventSource.addEventListener('logs_init', (e) => {
        logsContainer.innerHTML = '';
        JSON.parse(e.data).forEach(log =>
            logsContainer.appendChild(createLogElement(log)));
    });

    // Новые логи в реальном времени
    eventSource.addEventListener('logs_new', (e) => {
        JSON.parse(e.data).forEach(log =>
            logsContainer.appendChild(createLogElement(log, true)));
        if (autoscrollCb.checked) scrollToBottom();
    });
}'''
add_code_block(slide, Inches(0.6), Inches(1.65), Inches(7.5), Inches(5.3), code_js, font_size=10)

# Right: features
feature_tints = [BLUE_TINT, TEAL_TINT, ORANGE_TINT, RED_TINT]
features = [
    ("EventSource API", "Встроенный браузерный API\nАвтопереподключение при обрыве\nПоддержка Last-Event-Id", ACCENT_BLUE),
    ("Фильтрация", "Поиск устройств по имени\nЧекбоксы фильтрации по платформе\nДинамическое обновление", ACCENT_TEAL),
    ("Визуализация", "Цветовая подсветка уровней логов\nАнимация новых записей (CSS)\nAuto-scroll к последним записям", ACCENT_ORANGE),
    ("XSS-защита", "escapeHtml() для пользовательского\nввода — безопасное отображение", ACCENT_RED),
]

for i, (title, desc, color) in enumerate(features):
    cy = Inches(1.65 + i * 1.3)
    shape = add_shape(slide, Inches(8.4), cy, Inches(4.4), Inches(1.15), feature_tints[i], radius=0.03, border_color=BORDER_COLOR)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = 'Segoe UI'
    add_paragraph(tf, desc, size=11, color=TEXT_SECONDARY, space_before=Pt(4))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — DOWNLOAD & RECONNECT
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Скачивание логов и переподключение", "Chunked transfer + Last-Event-Id")

# Left
add_bullet_card(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.4),
                "GET /api/logs/download",
                ["Chunked Transfer-Encoding для потоковой отдачи",
                 "XRANGE пакетами по 5000 записей",
                 "Content-Disposition: attachment для скачивания",
                 "Не держит все данные в памяти"],
                title_color=ACCENT_ORANGE, tint_color=ORANGE_TINT)

add_bullet_card(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.4),
                "SSE Reconnect: Last-Event-Id",
                ["При обрыве — EventSource автоматически переподключается",
                 "Передаёт заголовок Last-Event-Id (Redis stream ID)",
                 "Сервер продолжает с этого ID через XREAD",
                 "Клиент не теряет ни одного сообщения"],
                title_color=ACCENT_PURPLE, tint_color=PURPLE_TINT)

# Code: download
code_dl = '''class logs_download_provider {
    bool operator()(size_t offset, httplib::DataSink& sink) {
        constexpr size_t batch_size = 5000;
        std::vector<stream_item> stream_data;

        redis_client_->xrange(stream_key_,
            start_id, "+", batch_size,
            std::back_inserter(stream_data));

        if (stream_data.empty()) {
            sink.done();
            return false;
        }

        std::string chunk;
        for (const auto& [id, fields] : stream_data)
            for (const auto& [key, value] : fields)
                if (key == "message") chunk += value;

        sink.write(chunk.c_str(), chunk.size());
        return stream_data.size() >= batch_size;
    }
};'''
add_code_block(slide, Inches(0.6), Inches(4.4), Inches(6.0), Inches(2.9), code_dl, font_size=10)

# Right: reconnect diagram
cy_base = Inches(4.5)
items_reconn = [
    ("Browser", "EventSource подключается", ACCENT_TEAL, TEAL_TINT),
    ("Server", "Отправляет id: 1679064645123-0", ACCENT_ORANGE, ORANGE_TINT),
    ("Обрыв!", "Соединение разорвано", ACCENT_RED, RED_TINT),
    ("Browser", "Reconnect + Last-Event-Id", ACCENT_TEAL, TEAL_TINT),
    ("Server", "XREAD с 1679064645123-0", ACCENT_ORANGE, ORANGE_TINT),
]
for i, (actor, desc, color, tint) in enumerate(items_reconn):
    cy = cy_base + Inches(i * 0.55)
    badge = add_shape(slide, Inches(7.0), cy, Inches(1.2), Inches(0.4), color, radius=0.1)
    set_text(badge, actor, size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(8.4), cy, Inches(4.3), Inches(0.4),
                 desc, size=11, color=TEXT_SECONDARY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — КЛЮЧЕВЫЕ РЕШЕНИЯ
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Ключевые проектные решения", "Архитектурные выборы и их обоснования")

decisions = [
    ("Redis Streams", "Упорядоченный лог с автогенерацией ID.\nВстроенная поддержка XRANGE/XREAD.\nЛучше подходит, чем LIST или PUBSUB.", ACCENT_BLUE, BLUE_TINT),
    ("Connection Pool", "redis::ConnectionPoolOptions — до 128\nсоединений. Каждый SSE-поток может\nчитать Redis без блокировки других.", ACCENT_TEAL, TEAL_TINT),
    ("Blocking XREAD", "XREAD с таймаутом заменяет\nbusy-polling. Минимальная нагрузка\nна CPU при ожидании новых записей.", ACCENT_ORANGE, ORANGE_TINT),
    ("condition_variable", "Когда лог не выбран — поток спит\nна CV вместо опроса Redis.\nМгновенное пробуждение при изменении.", ACCENT_PURPLE, PURPLE_TINT),
    ("Пакетная загрузка", "XRANGE по 1000 записей за раз.\nНе загружает весь стрим в память.\nSSE отправляет пакеты по мере чтения.", ACCENT_RED, RED_TINT),
    ("TTL 72 часа", "Все ключи автоматически удаляются.\nНе нужен фоновый сборщик мусора.\nRedis управляет памятью сам.", ACCENT_BLUE, BLUE_TINT),
]

for i, (title, desc, color, tint) in enumerate(decisions):
    col = i % 3
    row = i // 3
    cx = Inches(0.4 + col * 4.2)
    cy = Inches(1.7 + row * 2.7)

    shape = add_shape(slide, cx, cy, Inches(3.9), Inches(2.4), tint, radius=0.03, border_color=BORDER_COLOR)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = 'Segoe UI'

    add_paragraph(tf, desc, size=12, color=TEXT_SECONDARY, space_before=Pt(10))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — СТЕК ТЕХНОЛОГИЙ
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Стек технологий", "Библиотеки и стандарты")

tech_items = [
    ("C++23", "std::format, std::println, std::ranges,\nstd::jthread, std::stop_token,\nstd::chrono, ranges::to", ACCENT_BLUE, BLUE_TINT,
     "Язык"),
    ("redis-plus-plus", "Клиент Redis для C++.\nПоддержка Streams, Pub/Sub,\nConnection Pool, pipeline.", ACCENT_RED, RED_TINT,
     "Redis"),
    ("cpp-httplib", "Header-only HTTP/HTTPS сервер.\nContent provider для SSE.\nChunked transfer encoding.", ACCENT_ORANGE, ORANGE_TINT,
     "HTTP"),
    ("nlohmann/json", "JSON serialization/deserialization.\nИспользуется для API ответов\nи конфигурации.", ACCENT_TEAL, TEAL_TINT,
     "JSON"),
    ("EventSource", "Браузерный API для SSE.\nАвтоматическое переподключение.\nПоддержка Last-Event-Id.", ACCENT_PURPLE, PURPLE_TINT,
     "SSE"),
    ("HTML / CSS / JS", "Без фреймворков — vanilla JS.\nCSS custom properties, flexbox.\nАдаптивный тёмный лог-просмотрщик.", RGBColor(0xB4, 0x8B, 0x20), RGBColor(0xFE, 0xF9, 0xEF),
     "UI"),
]

for i, (name, desc, color, tint, category) in enumerate(tech_items):
    col = i % 3
    row = i // 3
    cx = Inches(0.4 + col * 4.2)
    cy = Inches(1.7 + row * 2.7)

    # Category badge
    badge = add_shape(slide, cx, cy, Inches(0.9), Inches(0.35), color, radius=0.15)
    set_text(badge, category, size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    shape = add_shape(slide, cx, cy + Inches(0.45), Inches(3.9), Inches(1.95), tint, radius=0.03, border_color=BORDER_COLOR)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = name
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = 'Segoe UI'

    add_paragraph(tf, desc, size=12, color=TEXT_SECONDARY, space_before=Pt(8))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — ИТОГИ
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

# Top/bottom accents
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_TEAL)
add_rect(slide, Inches(0), Inches(7.44), W, Inches(0.06), ACCENT_BLUE)

# Soft header bg
add_rect(slide, Inches(0), Inches(0.06), W, Inches(2.1), RGBColor(0xF0, 0xF7, 0xFF))

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.9),
             "Итоги", size=40, color=TEXT_PRIMARY, bold=True)
add_rect(slide, Inches(1), Inches(1.4), Inches(2.5), Inches(0.04), ACCENT_TEAL)

summary = [
    ("Клиент", "Простой C++23 клиент, отправляющий\nлоги через XADD в Redis Streams.\nРегистрация устройства + интерактивный режим.", ACCENT_BLUE, BLUE_TINT),
    ("Redis", "Центральное хранилище:\nHASH + SET + STREAM.\nKeyspace Notifications + TTL.", ACCENT_RED, RED_TINT),
    ("Сервер", "HTTP-сервер с SSE-стримингом.\nnotification_center + connection pool.\nПакетная загрузка + блокирующий XREAD.", ACCENT_ORANGE, ORANGE_TINT),
    ("Фронтенд", "Vanilla JS + EventSource.\nТри панели, фильтрация, поиск.\nЦветовая подсветка + auto-scroll.", ACCENT_TEAL, TEAL_TINT),
]

for i, (title, desc, color, tint) in enumerate(summary):
    cx = Inches(0.6 + i * 3.15)
    cy = Inches(2.5)
    shape = add_shape(slide, cx, cy, Inches(2.9), Inches(3.5), tint, radius=0.04, border_color=BORDER_COLOR)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # Color top bar on card
    add_rect(slide, cx, cy, Inches(2.9), Inches(0.06), color)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = 'Segoe UI'

    add_paragraph(tf, desc, size=12, color=TEXT_SECONDARY, space_before=Pt(12))

add_text_box(slide, Inches(1), Inches(6.4), Inches(11), Inches(0.6),
             "Полноценная система real-time логирования на C++23 без внешних веб-фреймворков",
             size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────

output_path = os.path.join(os.path.dirname(__file__), "presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
